from pathlib import Path
from pptx import Presentation

REPO_ROOT = Path(".")
OUTPUT_DIR = REPO_ROOT / "extracted-text"

OUTPUT_DIR.mkdir(exist_ok=True)

def extract_text_from_shape(shape):
    extracted = []

    if hasattr(shape, "text"):
        text = shape.text.strip()
        if text:
            extracted.append(text)

    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            row_values = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_values.append(cell_text)
            if any(row_values):
                extracted.append(" | ".join(row_values))

    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            extracted.extend(extract_text_from_shape(sub_shape))

    return extracted


def extract_pptx(pptx_path):
    prs = Presentation(pptx_path)

    lines = [
        f"# Extracted Text: {pptx_path.name}",
        "",
        f"Source file: `{pptx_path}`",
        ""
    ]

    for slide_number, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n## Slide {slide_number}\n")

        slide_text = []

        for shape in slide.shapes:
            slide_text.extend(extract_text_from_shape(shape))

        if slide_text:
            for text in slide_text:
                lines.append(text)
                lines.append("")
        else:
            lines.append("_No editable text found on this slide._")
            lines.append("")

        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append("### Speaker Notes")
                    lines.append("")
                    lines.append(notes_text)
                    lines.append("")
        except Exception:
            pass

    return "\n".join(lines)


def main():
    pptx_files = [
        path for path in REPO_ROOT.rglob("*.pptx")
        if ".git" not in path.parts
    ]

    if not pptx_files:
        print("No PPTX files found.")
        return

    for pptx_path in pptx_files:
        print(f"Extracting: {pptx_path}")

        markdown = extract_pptx(pptx_path)

        safe_name = pptx_path.stem.replace(" ", "-")
        output_path = OUTPUT_DIR / f"{safe_name}-extracted.md"

        output_path.write_text(markdown, encoding="utf-8")

        print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
